"""
S09 - Relatórios Mensais por Componente (Markdown + PPTX)
==========================================================
Etapa 9 do Pipeline (opcional)

O QUE FAZ:
Para cada equipamento, gera um relatório mensal que combina produção, qualidade
e a prescrição calculada em s08. Saídas:

- outputs/relatorios_mensais_componentes/EQ-XXX.md
- outputs/relatorios_mensais_componentes/INDEX.md
- outputs/relatorios_mensais_componentes_ppt/EQ-XXX.pptx
- outputs/relatorios_mensais_componentes_ppt/Apresentacao_Consolidada.pptx
- outputs/relatorios_mensais_componentes_ppt/relatorio_mensal_por_componente.zip

Também limpa quaisquer artefatos antigos com prefixo IJ-* nesses diretórios.

ENTRADAS:
- data/raw/EQ-*.csv
- outputs/historico_completo.csv         (s07)
- outputs/historico_recente.csv          (s07)
- outputs/janelas_operacao.csv           (s07)
- outputs/ociosidade.csv                 (s07)
- outputs/prescricao_manutencao.csv      (s08)
"""
from __future__ import annotations

import copy
import shutil
import sys
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = Path(__file__).parent
BASE_DIR = SCRIPT_DIR.parent
sys.path.insert(0, str(BASE_DIR / "config"))

try:
    from paths import (
        DATA_RAW_DIR,
        HISTORICO_COMPLETO_FILE, HISTORICO_RECENTE_FILE,
        JANELAS_OPERACAO_FILE, OCIOSIDADE_FILE, PRESCRICAO_MANUTENCAO_FILE,
        RELATORIOS_COMPONENTES_DIR, RELATORIOS_COMPONENTES_PPT_DIR,
    )
except ImportError:
    DATA_RAW_DIR = BASE_DIR / "data" / "raw"
    OUTPUTS_DIR = BASE_DIR / "outputs"
    HISTORICO_COMPLETO_FILE = OUTPUTS_DIR / "historico_completo.csv"
    HISTORICO_RECENTE_FILE = OUTPUTS_DIR / "historico_recente.csv"
    JANELAS_OPERACAO_FILE = OUTPUTS_DIR / "janelas_operacao.csv"
    OCIOSIDADE_FILE = OUTPUTS_DIR / "ociosidade.csv"
    PRESCRICAO_MANUTENCAO_FILE = OUTPUTS_DIR / "prescricao_manutencao.csv"
    RELATORIOS_COMPONENTES_DIR = OUTPUTS_DIR / "relatorios_mensais_componentes"
    RELATORIOS_COMPONENTES_PPT_DIR = OUTPUTS_DIR / "relatorios_mensais_componentes_ppt"


def _fmt_int(n) -> str:
    if pd.isna(n):
        return "—"
    return f"{int(n):,}".replace(",", ".")


def _fmt_pct(x) -> str:
    if pd.isna(x):
        return "—"
    return f"{x:,.2f} %".replace(",", "X").replace(".", ",").replace("X", ".")


def _fmt_dt(d) -> str:
    if pd.isna(d):
        return "—"
    return pd.Timestamp(d).strftime("%d/%m/%Y")


def _fmt_period_label(d) -> str:
    if pd.isna(d):
        return "—"
    months = ["Jan", "Fev", "Mar", "Abr", "Mai", "Jun",
              "Jul", "Ago", "Set", "Out", "Nov", "Dez"]
    d = pd.Timestamp(d)
    return f"{months[d.month - 1]}/{d.year}"


def _read_production() -> pd.DataFrame:
    files = sorted(DATA_RAW_DIR.glob("EQ-*.csv"))
    frames = []
    for f in files:
        df = pd.read_csv(f)
        df["arquivo_origem"] = f.name
        frames.append(df)
    if not frames:
        return pd.DataFrame()
    df = pd.concat(frames, ignore_index=True)
    df["data_producao"] = pd.to_datetime(df["Data de Produção"], dayfirst=True, errors="coerce")
    df["equipamento"] = df["Cód. Recurso"].astype(str).str.strip()
    df["ano_mes"] = df["data_producao"].dt.to_period("M")
    return df


def _build_md(equip: str, prod: pd.DataFrame, hist: pd.DataFrame,
              prescricao: pd.Series | None, hoje: pd.Timestamp) -> str:
    """Monta um relatório markdown completo para um equipamento."""
    grp = prod[prod["equipamento"] == equip].copy()
    if grp.empty:
        return f"# Relatório mensal por componente — {equip}\n\n_Sem apontamentos de produção registrados._\n"

    total_prod = int(grp["Qtd. Produzida"].fillna(0).sum())
    total_ref = int(grp["Qtd. Refugada"].fillna(0).sum())
    taxa_ref = (total_ref / total_prod * 100) if total_prod else 0.0
    inicio, fim = grp["data_producao"].min(), grp["data_producao"].max()
    n_compostos = grp["Descrição da massa (Composto)"].nunique()

    # Manutenções registradas
    mevents = hist[hist["equipamento"] == equip].copy() if not hist.empty else pd.DataFrame()
    mevents["data_evento"] = pd.to_datetime(mevents["data_evento"], errors="coerce") \
        if not mevents.empty else pd.NaT

    # Produção mensal por composto
    monthly = (
        grp.groupby(["ano_mes", "Descrição da massa (Composto)"])
           .agg(qtd=("Qtd. Produzida", "sum"),
                ref=("Qtd. Refugada", "sum"),
                dias=("data_producao", lambda s: s.dt.normalize().nunique()))
           .reset_index()
    )
    if not monthly.empty:
        monthly["pct_total"] = monthly["qtd"] / total_prod * 100

    # Resumo por composto
    resumo = (
        grp.groupby("Descrição da massa (Composto)")
           .agg(qtd=("Qtd. Produzida", "sum"),
                ref=("Qtd. Refugada", "sum"),
                meses=("ano_mes", "nunique"))
           .reset_index()
           .sort_values("qtd", ascending=False)
    )
    resumo["pct_total"] = resumo["qtd"] / total_prod * 100

    lines = [
        f"# Relatório mensal por componente — {equip}",
        "",
        f"**Fonte:** `data/raw/{equip}.csv`  ",
        f"**Período coberto:** {_fmt_dt(inicio)} → {_fmt_dt(fim)}  ",
        f"**Total produzido:** {_fmt_int(total_prod)} peças  ",
        f"**Total refugado:** {_fmt_int(total_ref)} peças ({_fmt_pct(taxa_ref)})  ",
        f"**Compostos distintos utilizados:** {n_compostos}",
        "",
        "## Manutenções registradas",
        "",
        "| Data | Evento | Dentro do período de produção? |",
        "|------|--------|-------------------------------|",
    ]
    if mevents.empty:
        lines.append("| — | _sem registros_ | — |")
    else:
        mevents = mevents.sort_values("data_evento")
        for _, r in mevents.iterrows():
            dentro = "✅ sim" if (pd.notna(r["data_evento"]) and inicio <= r["data_evento"] <= fim) else "—"
            tipo_label = {
                "ULTIMA_SUBSTITUICAO": "Última troca",
                "PENULTIMA_SUBSTITUICAO": "Penúltima troca",
                "PREVENTIVA": "Preventiva (#" + str(r.get("observacoes") or "?") + ")",
            }.get(r["tipo_evento"], r["tipo_evento"])
            lines.append(f"| {_fmt_dt(r['data_evento'])} | {tipo_label} | {dentro} |")

    lines += [
        "",
        "## Produção mensal por composto",
        "",
        "| Ano-Mês | Composto | Qtd. produzida | % do total | Refugo | Dias com produção |",
        "|---------|----------|---------------:|-----------:|-------:|------------------:|",
    ]
    for _, r in monthly.iterrows():
        lines.append(
            f"| {r['ano_mes']} | {r['Descrição da massa (Composto)']} | "
            f"{_fmt_int(r['qtd'])} | {_fmt_pct(r['pct_total'])} | "
            f"{_fmt_int(r['ref'])} | {int(r['dias'])} |"
        )

    lines += [
        "",
        "## Resumo por composto (período inteiro)",
        "",
        "| Composto | Qtd. produzida | % do total | Refugo | Meses com uso |",
        "|----------|---------------:|-----------:|-------:|--------------:|",
    ]
    for _, r in resumo.iterrows():
        lines.append(
            f"| {r['Descrição da massa (Composto)']} | {_fmt_int(r['qtd'])} | "
            f"{_fmt_pct(r['pct_total'])} | {_fmt_int(r['ref'])} | {int(r['meses'])} |"
        )
    lines.append(
        f"| **TOTAL** | **{_fmt_int(total_prod)}** | **100,00 %** | "
        f"**{_fmt_int(total_ref)}** | **{int(grp['ano_mes'].nunique())}** |"
    )

    if prescricao is not None and pd.notna(prescricao.get("data_prescrita")):
        lines += [
            "",
            "## Prescrição da próxima manutenção",
            "",
            f"**Data prescrita:** `{_fmt_dt(prescricao['data_prescrita'])}`  ",
            f"**Data da última troca:** `{_fmt_dt(prescricao.get('data_ultima_substituicao'))}`  ",
            f"**Data de referência (hoje):** `{_fmt_dt(hoje)}`  ",
            f"**Dias restantes:** `{int(prescricao['dias_restantes']) if pd.notna(prescricao['dias_restantes']) else '—'}`  ",
            f"**Urgência:** **{prescricao['urgencia']}**",
            "",
            "### Como foi calculada",
            "",
            "```text",
            "T_base          = mediana(dias_em_operacao) por equipamento (fallback 450)",
            "fator_desgaste  = clamp( 1 / (amp_atual / amp_mediana_hist) , 0.60 , 1.20 )",
            "fator_massa     = clamp( 1 / (massa_atual / massa_mediana_janelas) , 0.70 , 1.30 )",
            "T_prescrito     = T_base × fator_desgaste × fator_massa",
            "data_prescrita  = data_ultima_sub + T_prescrito + dias_ociosidade",
            "```",
            "",
            "### Valores usados para este equipamento",
            "",
            "| Campo | Valor |",
            "|-------|------:|",
            f"| T_base (dias) | {prescricao['T_base_dias']} |",
            f"| amplitude A atual | {prescricao.get('amplitude_a_atual')} |",
            f"| amplitude B atual | {prescricao.get('amplitude_b_atual')} |",
            f"| fator_desgaste | {prescricao['fator_desgaste']} |",
            f"| massa pós-última troca (kg) | {prescricao['massa_atual_kg']} |",
            f"| massa mediana janelas (kg) | {prescricao['massa_mediana_janelas_kg']} |",
            f"| fator_massa | {prescricao['fator_massa']} |",
            f"| T_prescrito (dias) | {prescricao['T_prescrito_dias']} |",
            f"| dias_ociosidade | {prescricao['dias_ociosidade']} |",
            "",
            "### Faixas de urgência",
            "",
            "- `dias_restantes < 0`  → **ATRASADO**",
            "- `0 ≤ dias_restantes < 30`  → **URGENTE**",
            "- `30 ≤ dias_restantes < 90`  → **ATENÇÃO**",
            "- `dias_restantes ≥ 90`  → **OK**",
        ]

    return "\n".join(lines) + "\n"


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_ACCENT = RGBColor(0xE2, 0x6E, 0x2C)
COLOR_LIGHT = RGBColor(0xF1, 0xF1, 0xF1)
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)


def _add_text(slide, left, top, width, height, text, *, size=14, bold=False,
              color=COLOR_TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _add_band(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


_CHART_COLORS = ["#1F3A5F", "#E26E2C", "#2E864E", "#E6B333", "#5B7DB1",
                 "#C0392B", "#888888"]


def _composto_summary(grp: pd.DataFrame, top_n: int = 6) -> pd.DataFrame:
    """Top N compostos por produção; agrega o resto em 'Outros'."""
    summary = (
        grp.groupby("Descrição da massa (Composto)")
           .agg(qtd=("Qtd. Produzida", "sum"),
                ref=("Qtd. Refugada", "sum"))
           .reset_index()
           .sort_values("qtd", ascending=False)
    )
    if len(summary) <= top_n:
        return summary
    head = summary.iloc[:top_n].copy()
    tail = summary.iloc[top_n:]
    outros = pd.DataFrame([{
        "Descrição da massa (Composto)": "Outros",
        "qtd": tail["qtd"].sum(),
        "ref": tail["ref"].sum(),
    }])
    return pd.concat([head, outros], ignore_index=True)


def _plot_pie_compostos(summary: pd.DataFrame, png_path: Path) -> Path | None:
    if summary.empty or summary["qtd"].sum() == 0:
        return None
    fig, ax = plt.subplots(figsize=(6.0, 4.0), dpi=150)
    ax.pie(
        summary["qtd"],
        labels=summary["Descrição da massa (Composto)"],
        autopct="%1.1f%%",
        colors=_CHART_COLORS[: len(summary)],
        textprops={"fontsize": 8},
        startangle=90,
    )
    ax.set_title("Distribuição da produção por composto", fontsize=11, color="#1F3A5F")
    fig.tight_layout()
    fig.savefig(png_path, bbox_inches="tight")
    plt.close(fig)
    return png_path


def _plot_bar_refugo(summary: pd.DataFrame, png_path: Path) -> Path | None:
    if summary.empty or summary["qtd"].sum() == 0:
        return None
    df = summary.copy()
    df["taxa_pct"] = (df["ref"] / df["qtd"].replace(0, pd.NA) * 100).fillna(0.0)
    fig, ax1 = plt.subplots(figsize=(7.5, 3.8), dpi=150)
    x = range(len(df))
    bars = ax1.bar(x, df["ref"], color="#E26E2C", label="Peças refugadas")
    ax1.set_ylabel("Peças refugadas", color="#E26E2C", fontsize=9)
    ax1.tick_params(axis="y", labelcolor="#E26E2C", labelsize=8)
    ax1.set_xticks(list(x))
    ax1.set_xticklabels(df["Descrição da massa (Composto)"], rotation=30, ha="right", fontsize=8)
    ax2 = ax1.twinx()
    ax2.plot(list(x), df["taxa_pct"], color="#1F3A5F", marker="o", linewidth=1.5, label="Taxa de refugo (%)")
    ax2.set_ylabel("Taxa de refugo (%)", color="#1F3A5F", fontsize=9)
    ax2.tick_params(axis="y", labelcolor="#1F3A5F", labelsize=8)
    for b, v in zip(bars, df["ref"]):
        ax1.text(b.get_x() + b.get_width() / 2, b.get_height(), _fmt_int(v),
                 ha="center", va="bottom", fontsize=7)
    fig.suptitle("Refugo por composto", fontsize=11, color="#1F3A5F")
    fig.tight_layout(rect=[0, 0, 1, 0.94])
    fig.savefig(png_path, bbox_inches="tight")
    plt.close(fig)
    return png_path


def _build_pptx(equip: str, prod: pd.DataFrame, hist: pd.DataFrame,
                prescricao: pd.Series | None, output_path: Path,
                hoje: pd.Timestamp) -> None:
    grp = prod[prod["equipamento"] == equip]
    total_prod = int(grp["Qtd. Produzida"].fillna(0).sum()) if not grp.empty else 0
    total_ref = int(grp["Qtd. Refugada"].fillna(0).sum()) if not grp.empty else 0
    taxa_ref = (total_ref / total_prod * 100) if total_prod else 0.0
    n_meses = grp["ano_mes"].nunique() if not grp.empty else 0
    n_compostos = grp["Descrição da massa (Composto)"].nunique() if not grp.empty else 0
    inicio = grp["data_producao"].min() if not grp.empty else pd.NaT
    fim = grp["data_producao"].max() if not grp.empty else pd.NaT
    period_label = f"{_fmt_period_label(inicio)} → {_fmt_period_label(fim)}"

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank = prs.slide_layouts[6]

    # --- Slide 1: capa ---
    s = prs.slides.add_slide(blank)
    _add_band(s, 0, 0, SLIDE_WIDTH_IN, 1.6, COLOR_PRIMARY)
    _add_text(s, 0.5, 0.3, 12, 0.6, equip, size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text(s, 0.5, 0.95, 12, 0.4, "Relatório Mensal por Componente",
              size=14, color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text(s, 0.5, 2.2, 12, 0.4, f"Análise de desempenho · {period_label}",
              size=16, bold=True, color=COLOR_PRIMARY)
    _add_text(s, 0.5, 2.7, 12, 0.4, "Produção · Qualidade · Manutenção",
              size=12, color=COLOR_TEXT)

    _add_text(s, 0.5, 3.8, 4, 0.5, _fmt_int(total_prod), size=36, bold=True, color=COLOR_PRIMARY)
    _add_text(s, 0.5, 4.5, 4, 0.4, "Peças produzidas", size=12)
    _add_text(s, 0.5, 4.85, 4, 0.4, period_label, size=10, color=COLOR_ACCENT)

    _add_text(s, 4.7, 3.8, 4, 0.5, str(n_compostos), size=36, bold=True, color=COLOR_PRIMARY)
    _add_text(s, 4.7, 4.5, 4, 0.4, "Compostos utilizados", size=12)
    _add_text(s, 4.7, 4.85, 4, 0.4, f"{n_meses} meses de operação", size=10, color=COLOR_ACCENT)

    _add_text(s, 8.9, 3.8, 4, 0.5, _fmt_pct(taxa_ref), size=36, bold=True, color=COLOR_ACCENT)
    _add_text(s, 8.9, 4.5, 4, 0.4, "Taxa de refugo", size=12)
    _add_text(s, 8.9, 4.85, 4, 0.4, f"{_fmt_int(total_ref)} peças refugadas", size=10, color=COLOR_TEXT)

    # --- Slide 2: visão executiva ---
    s = prs.slides.add_slide(blank)
    _add_band(s, 0, 0, SLIDE_WIDTH_IN, 0.8, COLOR_PRIMARY)
    _add_text(s, 0.4, 0.15, 11, 0.5, "Visão Executiva da Produção",
              size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text(s, 11.5, 0.2, 1.5, 0.4, equip, size=14, bold=True,
              color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.RIGHT)

    _add_text(s, 0.5, 1.2, 3, 0.5, _fmt_int(total_prod), size=28, bold=True, color=COLOR_PRIMARY)
    _add_text(s, 0.5, 1.85, 3, 0.4, "Total Produzido (peças)", size=11)
    _add_text(s, 3.7, 1.2, 3, 0.5, str(n_meses), size=28, bold=True, color=COLOR_PRIMARY)
    _add_text(s, 3.7, 1.85, 3, 0.4, "Meses de Operação", size=11)
    _add_text(s, 6.9, 1.2, 3, 0.5, str(n_compostos), size=28, bold=True, color=COLOR_PRIMARY)
    _add_text(s, 6.9, 1.85, 3, 0.4, "Compostos Distintos", size=11)
    _add_text(s, 10.0, 1.2, 3, 0.5, _fmt_pct(taxa_ref), size=28, bold=True, color=COLOR_ACCENT)
    _add_text(s, 10.0, 1.85, 3, 0.4, f"Taxa de Refugo ({_fmt_int(total_ref)} pç)", size=11)

    _add_band(s, 0.4, 2.6, 12.5, 4.4, COLOR_LIGHT)
    _add_text(s, 0.6, 2.7, 12, 0.4, "Destaques do Período", size=16, bold=True, color=COLOR_PRIMARY)
    if not grp.empty:
        monthly_total = grp.groupby("ano_mes")["Qtd. Produzida"].sum().sort_values(ascending=False)
        if not monthly_total.empty:
            top_mes = monthly_total.index[0]
            _add_text(s, 0.6, 3.3, 12, 0.4,
                      f"✓  Pico de produção em {top_mes}: {_fmt_int(monthly_total.iloc[0])} peças", size=12)
        composto_top = grp.groupby("Descrição da massa (Composto)")["Qtd. Produzida"].sum().sort_values(ascending=False)
        if not composto_top.empty:
            pct = composto_top.iloc[0] / total_prod * 100
            _add_text(s, 0.6, 3.8, 12, 0.4,
                      f"✓  {composto_top.index[0]} lidera com {_fmt_pct(pct)} da produção total", size=12)
        _add_text(s, 0.6, 4.3, 12, 0.4,
                  f"✓  Taxa de refugo de {_fmt_pct(taxa_ref)} no período", size=12)
        if prescricao is not None and pd.notna(prescricao.get("data_prescrita")):
            _add_text(s, 0.6, 4.8, 12, 0.4,
                      f"✓  Última troca em {_fmt_dt(prescricao.get('data_ultima_substituicao'))} — "
                      f"próxima prescrita em {_fmt_dt(prescricao['data_prescrita'])}",
                      size=12)

    # --- Slide 3: tendências mensais ---
    s = prs.slides.add_slide(blank)
    _add_band(s, 0, 0, SLIDE_WIDTH_IN, 0.8, COLOR_PRIMARY)
    _add_text(s, 0.4, 0.15, 11, 0.5, "Tendências de Produção Mensal",
              size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if not grp.empty:
        monthly_total = grp.groupby("ano_mes")["Qtd. Produzida"].sum()
        pico = int(monthly_total.max()) if not monthly_total.empty else 0
        media = int(monthly_total.mean()) if not monthly_total.empty else 0
        menor = int(monthly_total.min()) if not monthly_total.empty else 0
        _add_text(s, 0.5, 1.3, 4, 0.5, _fmt_int(pico), size=32, bold=True, color=COLOR_PRIMARY)
        _add_text(s, 0.5, 2.0, 4, 0.4, "Pico de Produção", size=12)
        _add_text(s, 4.7, 1.3, 4, 0.5, _fmt_int(media), size=32, bold=True, color=COLOR_PRIMARY)
        _add_text(s, 4.7, 2.0, 4, 0.4, "Média Mensal", size=12)
        _add_text(s, 8.9, 1.3, 4, 0.5, _fmt_int(menor), size=32, bold=True, color=COLOR_ACCENT)
        _add_text(s, 8.9, 2.0, 4, 0.4, "Menor Mês", size=12)

    # Diretório temporário pros PNGs deste equipamento
    tmpdir = Path(tempfile.mkdtemp(prefix=f"s09_{equip}_"))

    # --- Slide 4: compostos (tabela + pie chart) ---
    s = prs.slides.add_slide(blank)
    _add_band(s, 0, 0, SLIDE_WIDTH_IN, 0.8, COLOR_PRIMARY)
    _add_text(s, 0.4, 0.15, 11, 0.5, "Análise de Materiais e Compostos",
              size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if not grp.empty:
        composto_summary_full = (
            grp.groupby("Descrição da massa (Composto)")
               .agg(qtd=("Qtd. Produzida", "sum"),
                    ref=("Qtd. Refugada", "sum"),
                    meses=("ano_mes", "nunique"))
               .reset_index()
               .sort_values("qtd", ascending=False)
        )
        composto_summary = composto_summary_full.head(6)
        # Tabela à esquerda
        rows = len(composto_summary) + 1
        cols = 5
        tbl_shape = s.shapes.add_table(rows, cols, Inches(0.4), Inches(1.1),
                                        Inches(7.0), Inches(0.4 * rows))
        tbl = tbl_shape.table
        for i, h in enumerate(["Composto", "Produzido", "% Total", "Refugo", "Meses"]):
            cell = tbl.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(10)
        for ri, (_, r) in enumerate(composto_summary.iterrows(), start=1):
            tbl.cell(ri, 0).text = str(r["Descrição da massa (Composto)"])
            tbl.cell(ri, 1).text = _fmt_int(r["qtd"])
            tbl.cell(ri, 2).text = _fmt_pct(r["qtd"] / total_prod * 100)
            tbl.cell(ri, 3).text = _fmt_int(r["ref"])
            tbl.cell(ri, 4).text = str(int(r["meses"]))
            for ci in range(cols):
                for p in tbl.cell(ri, ci).text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(9)
        # Pie chart à direita
        pie_summary = _composto_summary(grp, top_n=6)
        pie_png = _plot_pie_compostos(pie_summary, tmpdir / "pie_compostos.png")
        if pie_png is not None:
            s.shapes.add_picture(str(pie_png), Inches(7.6), Inches(1.1),
                                 width=Inches(5.5), height=Inches(4.0))

    # --- Slide 5: qualidade (bar chart de refugo + métricas) ---
    s = prs.slides.add_slide(blank)
    _add_band(s, 0, 0, SLIDE_WIDTH_IN, 0.8, COLOR_PRIMARY)
    _add_text(s, 0.4, 0.15, 11, 0.5, "Controle de Qualidade — Taxa de Refugo",
              size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    # KPIs no topo
    _add_text(s, 0.5, 1.0, 4, 0.5, _fmt_pct(taxa_ref), size=32, bold=True, color=COLOR_ACCENT)
    _add_text(s, 0.5, 1.7, 4, 0.4, "Taxa de Refugo Global", size=11)
    _add_text(s, 4.7, 1.0, 4, 0.5, _fmt_int(total_ref), size=32, bold=True, color=COLOR_ACCENT)
    _add_text(s, 4.7, 1.7, 4, 0.4, "Peças Refugadas", size=11)
    _add_text(s, 8.9, 1.0, 4, 0.5, _fmt_int(total_prod), size=32, bold=True, color=COLOR_PRIMARY)
    _add_text(s, 8.9, 1.7, 4, 0.4, "Peças Produzidas", size=11)
    nivel = "Dentro do limite aceitável" if taxa_ref < 1.5 else "Acima da meta — investigar"
    _add_text(s, 0.5, 2.3, 12, 0.4, nivel, size=12, bold=True,
              color=COLOR_PRIMARY if taxa_ref < 1.5 else COLOR_ACCENT)
    # Bar chart de refugo por composto
    if not grp.empty:
        bar_summary = _composto_summary(grp, top_n=6)
        bar_png = _plot_bar_refugo(bar_summary, tmpdir / "bar_refugo.png")
        if bar_png is not None:
            s.shapes.add_picture(str(bar_png), Inches(0.5), Inches(2.9),
                                 width=Inches(12.3), height=Inches(4.3))

    # --- Slide 6: manutenção (prescrição com tabela 3 colunas) ---
    s = prs.slides.add_slide(blank)
    _add_band(s, 0, 0, SLIDE_WIDTH_IN, 0.8, COLOR_PRIMARY)
    _add_text(s, 0.4, 0.15, 11, 0.5, "Estratégia de Manutenção Preventiva",
              size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if prescricao is not None and pd.notna(prescricao.get("data_prescrita")):
        urg = prescricao["urgencia"]
        urg_color = {
            "ATRASADO": RGBColor(0xC0, 0x39, 0x2B),
            "URGENTE": COLOR_ACCENT,
            "ATENÇÃO": RGBColor(0xE6, 0xB3, 0x33),
            "OK": RGBColor(0x2E, 0x86, 0x4E),
        }.get(urg, COLOR_TEXT)
        _add_text(s, 0.5, 1.2, 5, 0.5, urg, size=28, bold=True, color=urg_color)
        _add_text(s, 0.5, 1.9, 5, 0.5,
                  f"{int(prescricao['dias_restantes']) if pd.notna(prescricao['dias_restantes']) else '—'} dias restantes",
                  size=18)
        _add_text(s, 0.5, 2.6, 8, 0.5, f"Próxima manutenção: {_fmt_dt(prescricao['data_prescrita'])}",
                  size=14, bold=True, color=COLOR_PRIMARY)
        _add_text(s, 0.5, 3.1, 8, 0.4, f"Última troca: {_fmt_dt(prescricao.get('data_ultima_substituicao'))}",
                  size=12)
        # Tabela 3 colunas: Campo | Valor | Efeito
        items = [
            ("T_base (dias)", str(prescricao["T_base_dias"]),
             "Mediana histórica do equipamento (fallback 450)"),
            ("Fator desgaste", str(prescricao["fator_desgaste"]),
             "<1 acelera (componente desgastado), >1 estende"),
            ("Fator massa", str(prescricao["fator_massa"]),
             "<1 acelera (massa pesada), >1 estende"),
            ("T_prescrito (dias)", str(prescricao["T_prescrito_dias"]),
             "T_base × fator_desgaste × fator_massa"),
            ("Dias ociosidade", str(prescricao["dias_ociosidade"]),
             "Soma à data prescrita (compensa dias parado)"),
        ]
        rows = len(items) + 1
        tbl_shape = s.shapes.add_table(rows, 3, Inches(0.5), Inches(4.0),
                                        Inches(12.3), Inches(0.45 * rows))
        tbl = tbl_shape.table
        for i, h in enumerate(["Campo", "Valor", "Efeito"]):
            cell = tbl.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(11)
        for ri, (k, v, efeito) in enumerate(items, start=1):
            tbl.cell(ri, 0).text = k
            tbl.cell(ri, 1).text = v
            tbl.cell(ri, 2).text = efeito
            for ci in range(3):
                for p in tbl.cell(ri, ci).text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(10)

    # --- Slide 7: conclusões ---
    s = prs.slides.add_slide(blank)
    _add_band(s, 0, 0, SLIDE_WIDTH_IN, 0.8, COLOR_PRIMARY)
    _add_text(s, 0.4, 0.15, 11, 0.5, "Conclusões e Insights Operacionais",
              size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text(s, 0.5, 1.2, 12, 0.5, "Principais Conclusões", size=16, bold=True, color=COLOR_PRIMARY)
    bullets = [
        f"Produção total de {_fmt_int(total_prod)} peças em {n_meses} mês(es)",
        f"Taxa de refugo de {_fmt_pct(taxa_ref)} — {_fmt_int(total_ref)} peças refugadas",
        f"{n_compostos} compostos distintos utilizados no período",
    ]
    if prescricao is not None and pd.notna(prescricao.get("data_prescrita")):
        bullets.append(
            f"Próxima manutenção prescrita: {_fmt_dt(prescricao['data_prescrita'])} "
            f"({int(prescricao['dias_restantes']) if pd.notna(prescricao['dias_restantes']) else '—'} dias restantes)"
        )
    for i, b in enumerate(bullets):
        _add_text(s, 0.7, 1.9 + i * 0.6, 12, 0.5, "• " + b, size=12)

    prs.save(str(output_path))
    shutil.rmtree(tmpdir, ignore_errors=True)


def _merge_pptx(pptx_files: list[Path], output_path: Path) -> None:
    """Concatena os PPTX gerados em uma única apresentação."""
    if not pptx_files:
        return
    base = Presentation(str(pptx_files[0]))
    for f in pptx_files[1:]:
        src = Presentation(str(f))
        for slide in src.slides:
            blank = base.slide_layouts[6]
            new = base.slides.add_slide(blank)
            for shape in slide.shapes:
                el = shape.element
                new.shapes._spTree.insert_element_before(copy.deepcopy(el), "p:extLst")
    base.save(str(output_path))


def _zip_dir(src_dir: Path, zip_path: Path) -> None:
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in sorted(src_dir.iterdir()):
            if p.is_file() and not p.name.endswith(".zip"):
                zf.write(p, arcname=p.name)


def _purge_legacy(dirs: list[Path]) -> None:
    for d in dirs:
        if not d.exists():
            continue
        for p in d.iterdir():
            if p.is_file() and (p.name.startswith("IJ-") or "IJ-" in p.name or p.name.startswith("Apresentação")):
                p.unlink()


def main(**kwargs) -> dict:
    print("=" * 60)
    print("ETAPA 9: RELATÓRIOS MENSAIS POR COMPONENTE")
    print("=" * 60)

    RELATORIOS_COMPONENTES_DIR.mkdir(parents=True, exist_ok=True)
    RELATORIOS_COMPONENTES_PPT_DIR.mkdir(parents=True, exist_ok=True)

    print("\n[0/4] Limpando artefatos legados (IJ-*)...")
    _purge_legacy([RELATORIOS_COMPONENTES_DIR, RELATORIOS_COMPONENTES_PPT_DIR])

    print("\n[1/4] Carregando insumos...")
    prod = _read_production()
    hist = pd.read_csv(HISTORICO_COMPLETO_FILE) if HISTORICO_COMPLETO_FILE.exists() else pd.DataFrame()
    if not hist.empty:
        hist["data_evento"] = pd.to_datetime(hist["data_evento"], errors="coerce")
    prescricao = pd.read_csv(PRESCRICAO_MANUTENCAO_FILE, parse_dates=[
        "data_ultima_substituicao", "data_prescrita", "data_referencia",
    ]) if PRESCRICAO_MANUTENCAO_FILE.exists() else pd.DataFrame()
    print(f"  ✓ {prod['equipamento'].nunique() if not prod.empty else 0} equipamentos com produção")

    hoje = kwargs.get("data_referencia") or datetime.now()
    if isinstance(hoje, str):
        hoje = pd.to_datetime(hoje, dayfirst=True)
    hoje = pd.Timestamp(hoje).normalize()

    equipamentos = sorted(prod["equipamento"].dropna().unique()) if not prod.empty else []

    print(f"\n[2/4] Gerando .md e .pptx para {len(equipamentos)} equipamentos...")
    pptx_files = []
    index_rows = []
    for equip in equipamentos:
        prescr = None
        if not prescricao.empty:
            sel = prescricao[prescricao["equipamento"] == equip]
            if not sel.empty:
                prescr = sel.iloc[0]

        md_text = _build_md(equip, prod, hist, prescr, hoje)
        md_path = RELATORIOS_COMPONENTES_DIR / f"{equip}.md"
        md_path.write_text(md_text, encoding="utf-8")

        pptx_path = RELATORIOS_COMPONENTES_PPT_DIR / f"{equip}.pptx"
        _build_pptx(equip, prod, hist, prescr, pptx_path, hoje)
        pptx_files.append(pptx_path)

        n_apont = int((prod["equipamento"] == equip).sum())
        index_rows.append((equip, n_apont))

    print(f"  ✓ {len(pptx_files)} pares md/pptx gerados")

    print("\n[3/4] INDEX.md e apresentação consolidada...")
    index_lines = [
        "# Índice — relatórios mensais por componente",
        "",
        f"_Gerado em {hoje.strftime('%d/%m/%Y')}._  ",
        f"_Equipamentos: {len(index_rows)}._",
        "",
        "| Equipamento | Apontamentos | Arquivo |",
        "|-------------|-------------:|---------|",
    ]
    for equip, n in index_rows:
        index_lines.append(f"| {equip} | {_fmt_int(n)} | [{equip}.md]({equip}.md) |")
    (RELATORIOS_COMPONENTES_DIR / "INDEX.md").write_text("\n".join(index_lines) + "\n", encoding="utf-8")

    consolidated = RELATORIOS_COMPONENTES_PPT_DIR / "Apresentacao_Consolidada.pptx"
    if pptx_files:
        _merge_pptx(pptx_files, consolidated)
        print(f"  ✓ Apresentação consolidada: {consolidated.name}")

    print("\n[4/4] Gerando ZIP...")
    zip_path = RELATORIOS_COMPONENTES_PPT_DIR / "relatorio_mensal_por_componente.zip"
    _zip_dir(RELATORIOS_COMPONENTES_PPT_DIR, zip_path)
    print(f"  ✓ {zip_path.name}")

    md_files = sorted(RELATORIOS_COMPONENTES_DIR.glob("EQ-*.md"))
    pdf_files = sorted(RELATORIOS_COMPONENTES_PPT_DIR.glob("EQ-*.pdf"))
    total_apontamentos = sum(n for _, n in index_rows)
    zip_kb = round(zip_path.stat().st_size / 1024, 1) if zip_path.exists() else 0
    consolidado_kb = round(consolidated.stat().st_size / 1024, 1) if pptx_files and consolidated.exists() else 0

    print("\n" + "=" * 60)
    print("ETAPA 9 CONCLUÍDA")
    print("=" * 60)
    return {
        "status": "success",
        "n_equipamentos": len(equipamentos),
        "n_md": len(md_files),
        "n_pptx": len(pptx_files),
        "n_pdf": len(pdf_files),
        "total_apontamentos": total_apontamentos,
        "consolidado": str(consolidated) if pptx_files else None,
        "consolidado_kb": consolidado_kb,
        "zip": str(zip_path),
        "zip_kb": zip_kb,
        "data_referencia": hoje.strftime("%Y-%m-%d"),
    }


if __name__ == "__main__":
    main()
